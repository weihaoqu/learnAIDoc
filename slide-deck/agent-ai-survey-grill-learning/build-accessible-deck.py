#!/usr/bin/env python3
"""Build the Agent AI teaching deck from one structured, accessible source.

Outputs:
  - native-text PPTX (editable text and vector shapes)
  - semantic HTML (screen-reader and browser fallback)
  - tagged PDF/UA-1 generated from the semantic HTML
  - contact sheet rendered from the PDF

The legacy SVG/PNG renderer is kept for provenance, but this script is the
canonical build path for classroom-facing artifacts.
"""

from __future__ import annotations

import html
import math
import shutil
import subprocess
from pathlib import Path
from typing import Any

import fitz
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
PUBLIC = ROOT.parents[1] / "assets" / "decks" / ROOT.name
SLUG = "agent-ai-survey-grill-learning"
PAPER_URL = "https://arxiv.org/abs/2401.03568"

COLORS = {
    "bg": "F5F0E6",
    "paper": "FFFCF4",
    "ink": "171717",
    "muted": "5E6670",
    "line": "D8D0C0",
    "teal": "2F7373",
    "maroon": "722F37",
    "brown": "8B7355",
    "cobalt": "2563EB",
    "amber": "B7791F",
    "red": "B42318",
    "green": "16803C",
    "violet": "6D28D9",
}


def card(title: str, body: str, color: str = "teal") -> dict[str, str]:
    return {"title": title, "body": body, "color": color}


SLIDES: list[dict[str, Any]] = [
    {
        "stage": "Agent AI survey",
        "title": "Agent AI Is Not Just a Chatbot",
        "subtitle": "A classroom reading of Agent AI: Surveying the Horizons of Multimodal Interaction",
        "bullets": [
            "Perception: what can the system observe?",
            "Action: what can it change?",
            "Feedback: what does it learn from next?",
        ],
        "layout": "flow",
        "cards": [
            card("Perception", "Reads language, images, audio, video, or environment state.", "cobalt"),
            card("Action", "Changes a physical or digital environment.", "maroon"),
            card("Feedback", "Observes the result and adjusts the next step.", "green"),
        ],
        "callout": "Agent AI participates in a task loop; it does more than return advice.",
    },
    {
        "stage": "2024 paper framing",
        "title": "The Paper Gives a Map, Not a Finished Machine",
        "subtitle": "Use the January 2024 survey as conceptual vocabulary—not as a 2026 catalog of current systems.",
        "bullets": [
            "Moves beyond the beginner formula “LLM + tools + memory.”",
            "Centers multimodal perception and grounded environments.",
            "Connects action prediction, feedback, evaluation, and risk.",
        ],
        "layout": "grid",
        "cards": [
            card("Foundation model", "A building block, not the entire agent.", "brown"),
            card("Multimodal perception", "Language plus sensory and contextual signals.", "cobalt"),
            card("Grounded action", "Behavior occurs inside a physical or virtual environment.", "maroon"),
            card("Evaluation in context", "Measure outcomes, feedback, and risk—not only text quality.", "teal"),
        ],
        "callout": "Historical framing: durable concepts, dated examples.",
    },
    {
        "stage": "Concept correction",
        "title": "Multimodal Is Not Automatically Agentic",
        "subtitle": "The difference is the task loop—not simply the number of media types a model can read.",
        "bullets": [
            "A model can understand images or video without choosing an action.",
            "An agent uses perception, state, goals, action, and feedback.",
            "A caption is useful output, but it is not yet an agent loop.",
        ],
        "layout": "compare",
        "cards": [
            card("Multimodal", "Understands different information sources: text, image, video, audio, or screen state.", "cobalt"),
            card("Agentic", "Uses perception plus state, goal, action, and feedback to decide what happens next.", "teal"),
        ],
        "callout": "More media is not the same as more agency.",
    },
    {
        "stage": "Hypothetical classroom test",
        "title": "Description Is Not Yet an Agent Loop",
        "subtitle": "Treat “students look confused” as an uncertain interpretation—not as a reliable fact about students.",
        "bullets": [
            "Describe a video signal: multimodal perception.",
            "Choose a reteaching path, quiz, observe, and adapt: agentic loop.",
            "The interpretation can be wrong before any action occurs.",
        ],
        "layout": "flow",
        "cards": [
            card("Perceive", "A system processes a hypothetical classroom signal.", "cobalt"),
            card("Interpret", "It predicts possible confusion—with uncertainty.", "amber"),
            card("Act + adapt", "It recommends a response and observes feedback.", "teal"),
        ],
        "callout": "Class activity only: do not record or upload real student audio or video.",
    },
    {
        "stage": "Reusable agent test",
        "title": "The Agent Loop Has Six Questions",
        "subtitle": "Use the same checklist for a robot, game agent, classroom tool, or software assistant.",
        "bullets": [
            "Trace what enters the system and what it can change.",
            "Name the feedback signal and the success evidence.",
            "Evaluate risk before granting authority.",
        ],
        "layout": "six",
        "cards": [
            card("1 · Perception", "What can it observe?", "cobalt"),
            card("2 · Environment", "Where is it grounded?", "brown"),
            card("3 · Action", "What can it change?", "maroon"),
            card("4 · Feedback", "What happens next?", "green"),
            card("5 · Evaluation", "What proves progress?", "teal"),
            card("6 · Ethical risk", "Who could be harmed?", "red"),
        ],
        "callout": "Capability and authority are separate design decisions.",
    },
    {
        "stage": "Ethics stakes",
        "title": "Wrong Perception Can Become Wrong Action",
        "subtitle": "Signals such as posture, accent, eye movement, lighting, or camera quality are ambiguous—not evidence of intent.",
        "bullets": [
            "Privacy: screens, faces, voice, grades, and teacher materials.",
            "Fairness: disability, accent, environment, and unequal technology.",
            "Accountability: who authorized the action and who corrects it?",
        ],
        "layout": "flow",
        "cards": [
            card("Sensitive signal", "Screen, face, voice, grade, or classroom behavior.", "cobalt"),
            card("AI interpretation", "A probabilistic and potentially biased inference.", "amber"),
            card("Consequence", "Reteach, flag, report, restrict, or change a learning path.", "red"),
        ],
        "callout": "Higher agency raises the cost of a mistaken inference.",
    },
    {
        "stage": "Learning boundary",
        "title": "AI Helps Learning When Students Still Own the Thinking",
        "subtitle": "The boundary should be observable student understanding—not merely polished final output.",
        "bullets": [
            "Helpful support can explain, question, plan, revise, and verify.",
            "Replacement produces the main answer without understanding.",
            "Students should be able to explain, verify, adapt, and defend.",
        ],
        "layout": "compare",
        "cards": [
            card("Helps learning", "The student makes decisions, checks evidence, and can defend the work.", "green"),
            card("Replaces learning", "The system supplies the reasoning while the student submits the surface output.", "red"),
        ],
        "callout": "Ask: can the student explain, verify, adapt, and defend the result?",
    },
    {
        "stage": "Operational policy",
        "title": "“AI Allowed” and “AI Banned” Are Both Too Weak",
        "subtitle": "A usable policy translates ethical language into observable permissions, duties, and boundaries.",
        "bullets": [
            "Name approved tools and allowed learning uses.",
            "Name forbidden uses and exam boundaries.",
            "Require privacy-safe disclosure and human verification.",
        ],
        "layout": "grid",
        "cards": [
            card("Tools", "Which systems are approved?", "cobalt"),
            card("Uses", "What help is permitted?", "green"),
            card("Prohibitions", "What may not be delegated?", "red"),
            card("Disclosure", "What assistance must be reported?", "teal"),
        ],
        "callout": "Operational policy replaces slogans with responsibilities.",
    },
    {
        "stage": "Classroom governance",
        "title": "Boundaries Must Come Before Deployment",
        "subtitle": "These are discussion principles for classroom-agent design—not a replacement for institutional policy.",
        "bullets": [
            "Permission before recording or analyzing people.",
            "Logs and audits for recommendations and actions.",
            "Human authority for grades, discipline, accommodations, and appeals.",
            "No unauthorized upload of course or peer materials.",
        ],
        "layout": "grid",
        "cards": [
            card("Permission", "People and data are not automatically available to an agent.", "cobalt"),
            card("Audit", "Recommendations and actions need a reviewable trail.", "teal"),
            card("Authority", "High-stakes decisions remain with accountable humans.", "maroon"),
            card("Materials", "Protect course, peer, and restricted information.", "red"),
        ],
        "callout": "An agent may suggest; accountable humans must decide and hear appeals.",
    },
    {
        "stage": "Takeaway",
        "title": "From Advice to an Accountable Task Loop",
        "subtitle": "The upgrade is AI participating in perception, action, and feedback—inside explicit human boundaries.",
        "bullets": [
            "Chatbot: returns advice or content.",
            "Agent AI: participates in the task loop.",
            "Governance: who authorized, who checks, and who is accountable?",
        ],
        "layout": "flow",
        "cards": [
            card("Advice", "A chatbot proposes what a person might do.", "brown"),
            card("Task loop", "An agent perceives, acts, and receives feedback.", "teal"),
            card("Accountability gate", "A human reviews authority, evidence, risk, and consequences.", "red"),
        ],
        "callout": "Design the authority boundary—not only the capability.",
    },
]


def rgb(name: str) -> RGBColor:
    value = COLORS.get(name, name).lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def tint(name: str, ratio: float = 0.90) -> RGBColor:
    source = rgb(name)
    return RGBColor(
        round(source[0] + (255 - source[0]) * ratio),
        round(source[1] + (255 - source[1]) * ratio),
        round(source[2] + (255 - source[2]) * ratio),
    )


def set_shape_name(shape: Any, name: str, description: str | None = None) -> None:
    nodes = shape._element.xpath(".//p:cNvPr")
    if not nodes:
        return
    nodes[0].set("name", name)
    if description:
        nodes[0].set("descr", description)


def add_box(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = "paper",
    line: str = "line",
    radius: bool = True,
    name: str = "Panel",
) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.2)
    set_shape_name(shape, name)
    return shape


def add_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 20,
    color: str = "ink",
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    name: str = "Text",
    margin: float = 0.02,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_name(shape, name)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_bullets(slide: Any, items: list[str], x: float, y: float, w: float, h: float) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_name(shape, "Key teaching points", "Key teaching points for this slide")
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {item}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = rgb("ink")
        paragraph.space_after = Pt(11)
        paragraph.line_spacing = 1.05
    return shape


def add_arrow(slide: Any, x: float, y: float, w: float = 0.30) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("teal")
    shape.line.fill.background()
    set_shape_name(shape, "Flow arrow")


def add_card_shape(slide: Any, item: dict[str, str], x: float, y: float, w: float, h: float, index: int) -> None:
    color = item["color"]
    add_box(slide, x, y, w, h, fill=COLORS[color], line=color, name=f"Concept card {index}: {item['title']}")
    # Make the solid fill a very light tint while keeping the card native.
    shape = slide.shapes[-1]
    shape.fill.fore_color.rgb = tint(color, 0.92)
    add_text(
        slide,
        item["title"],
        x + 0.16,
        y + 0.17,
        w - 0.32,
        0.42,
        size=18 if w >= 2.6 else 15,
        color=color,
        bold=True,
        name=f"Card {index} title",
    )
    add_text(
        slide,
        item["body"],
        x + 0.16,
        y + 0.66,
        w - 0.32,
        h - 0.80,
        size=12.5 if w >= 2.6 else 10.5,
        color="ink",
        name=f"Card {index} explanation",
    )


def add_visual(slide: Any, spec: dict[str, Any]) -> None:
    cards = spec["cards"]
    layout = spec["layout"]
    x0, y0, width, height = 5.35, 2.82, 7.10, 3.58
    if layout == "flow":
        gap = 0.38
        card_w = (width - 2 * gap) / 3
        for index, item in enumerate(cards):
            x = x0 + index * (card_w + gap)
            add_card_shape(slide, item, x, y0 + 0.40, card_w, 2.45, index + 1)
            if index < len(cards) - 1:
                add_arrow(slide, x + card_w + 0.05, y0 + 1.40, 0.25)
    elif layout == "compare":
        gap = 0.40
        card_w = (width - gap) / 2
        for index, item in enumerate(cards):
            add_card_shape(slide, item, x0 + index * (card_w + gap), y0 + 0.28, card_w, 2.75, index + 1)
    elif layout == "six":
        gap_x, gap_y = 0.22, 0.22
        card_w = (width - 2 * gap_x) / 3
        card_h = 1.50
        for index, item in enumerate(cards):
            row, col = divmod(index, 3)
            add_card_shape(
                slide,
                item,
                x0 + col * (card_w + gap_x),
                y0 + row * (card_h + gap_y),
                card_w,
                card_h,
                index + 1,
            )
    else:
        gap_x, gap_y = 0.32, 0.24
        card_w = (width - gap_x) / 2
        card_h = 1.52
        for index, item in enumerate(cards):
            row, col = divmod(index, 2)
            add_card_shape(
                slide,
                item,
                x0 + col * (card_w + gap_x),
                y0 + row * (card_h + gap_y),
                card_w,
                card_h,
                index + 1,
            )


def add_footer(slide: Any, number: int) -> None:
    source = add_text(
        slide,
        "Durante et al. · Agent AI: Surveying the Horizons of Multimodal Interaction · arXiv:2401.03568 · January 2024",
        0.72,
        7.16,
        10.9,
        0.20,
        size=8.5,
        color="muted",
        name="Source citation",
    )
    run = source.text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = PAPER_URL
    add_text(slide, f"{number} / {len(SLIDES)}", 11.72, 7.14, 0.62, 0.22, size=8.5, color="muted", align=PP_ALIGN.RIGHT, name="Slide number")


def build_pptx(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "Agent AI Survey — Multimodal and Embodied Agents Beyond Chatbots"
    presentation.core_properties.subject = "LearnAI classroom deck on multimodal and agentic systems"
    presentation.core_properties.author = "LearnAI"
    presentation.core_properties.keywords = "Agent AI, multimodal AI, embodied agents, AI ethics, education"
    presentation.core_properties.comments = "Native-text accessible rebuild generated from build-accessible-deck.py"

    for number, spec in enumerate(SLIDES, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = rgb("bg")
        add_box(slide, 0.28, 0.28, 12.77, 6.78, fill="paper", line="line", name="Decorative slide panel")
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.28), Inches(1.17), Inches(12.77), Inches(0.012))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb("line")
        line.line.fill.background()
        set_shape_name(line, "Decorative header divider")

        badge_w = min(2.95, max(1.45, 0.095 * len(spec["stage"]) + 0.60))
        badge = add_box(slide, 0.70, 0.58, badge_w, 0.37, fill="paper", line="teal", name="Section label")
        badge.fill.fore_color.rgb = tint("teal", 0.88)
        add_text(slide, spec["stage"].upper(), 0.86, 0.655, badge_w - 0.30, 0.20, size=10.5, color="teal", bold=True, name="Section label text")
        add_text(slide, spec["title"], 0.70, 1.38, 11.85, 0.60, size=30, color="ink", bold=True, name="Slide title")
        add_text(slide, spec["subtitle"], 0.73, 2.12, 11.55, 0.48, size=14.5, color="muted", name="Slide subtitle")
        add_bullets(slide, spec["bullets"], 0.77, 3.05, 4.18, 3.10)
        add_visual(slide, spec)

        callout = add_box(slide, 5.55, 6.45, 6.66, 0.43, fill="paper", line="teal", name="Takeaway callout")
        callout.fill.fore_color.rgb = tint("teal", 0.90)
        add_text(slide, spec["callout"], 5.75, 6.55, 6.26, 0.22, size=11.5, color="teal", bold=True, align=PP_ALIGN.CENTER, name="Takeaway text")
        add_footer(slide, number)

    presentation.save(path)


def html_card(item: dict[str, str]) -> str:
    return (
        f'<article class="concept {html.escape(item["color"])}">'
        f'<h2>{html.escape(item["title"])}</h2>'
        f'<p>{html.escape(item["body"])}</p>'
        "</article>"
    )


def build_html(path: Path) -> None:
    slide_html: list[str] = []
    for number, spec in enumerate(SLIDES, start=1):
        items = "".join(f"<li>{html.escape(item)}</li>" for item in spec["bullets"])
        cards = "".join(html_card(item) for item in spec["cards"])
        slide_html.append(
            f'''<section class="slide" id="slide-{number}" aria-labelledby="title-{number}">
  <div class="panel">
    <p class="stage">{html.escape(spec["stage"])}</p>
    <h1 id="title-{number}">{html.escape(spec["title"])}</h1>
    <p class="subtitle">{html.escape(spec["subtitle"])}</p>
    <div class="content">
      <div class="teaching-points"><h2>Key teaching points</h2><ul>{items}</ul></div>
      <div class="visual {html.escape(spec["layout"])}" aria-label="Concept diagram">{cards}</div>
    </div>
    <p class="callout">{html.escape(spec["callout"])}</p>
    <footer><a href="{PAPER_URL}">Durante et al., Agent AI: Surveying the Horizons of Multimodal Interaction, arXiv:2401.03568 (January 2024)</a><span>{number} / {len(SLIDES)}</span></footer>
  </div>
</section>'''
        )

    document = f'''<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width,initial-scale=1">
  <meta name="description" content="Accessible LearnAI teaching deck about multimodal and embodied Agent AI.">
  <title>Agent AI Survey — Multimodal and Embodied Agents Beyond Chatbots</title>
  <style>
  @page {{ size: 13.333in 7.5in; margin: 0; }}
  :root {{ --bg:#F5F0E6;--paper:#FFFCF4;--ink:#171717;--muted:#5E6670;--line:#D8D0C0;--teal:#2F7373;--maroon:#722F37;--brown:#8B7355;--cobalt:#2563EB;--amber:#B7791F;--red:#B42318;--green:#16803C; }}
  * {{ box-sizing:border-box; }}
  html,body {{ margin:0;background:var(--bg);color:var(--ink);font-family:Arial,"Helvetica Neue",sans-serif; }}
  .slide {{ width:13.333in;height:7.5in;page-break-after:always;padding:.28in;background:var(--bg); }}
  .slide:last-child {{ page-break-after:auto; }}
  .panel {{ position:relative;width:100%;height:100%;padding:.30in .42in .22in;border:1px solid var(--line);border-radius:.08in;background:rgba(255,252,244,.96);display:flex;flex-direction:column; }}
  .stage {{ align-self:flex-start;margin:0 0 .12in;padding:.07in .15in;border:1px solid var(--teal);border-radius:.06in;background:#E6F0EC;color:var(--teal);font-size:10pt;font-weight:700;letter-spacing:.04em;text-transform:uppercase; }}
  h1 {{ margin:.07in 0 0;font-size:30pt;line-height:1.03; }}
  .subtitle {{ margin:.12in 0 .15in;color:var(--muted);font-size:14pt;line-height:1.22; }}
  .content {{ flex:1;display:grid;grid-template-columns:34% 66%;gap:.24in;min-height:0; }}
  .teaching-points {{ padding:.10in .08in; }}
  .teaching-points h2 {{ margin:0 0 .12in;color:var(--teal);font-size:12pt;text-transform:uppercase;letter-spacing:.04em; }}
  .teaching-points ul {{ margin:0;padding-left:.24in;font-size:15pt;line-height:1.28; }}
  .teaching-points li {{ margin:0 0 .16in; }}
  .visual {{ display:grid;gap:.15in;align-content:center;min-height:0; }}
  .visual.flow {{ grid-template-columns:repeat(3,1fr); }}
  .visual.compare {{ grid-template-columns:repeat(2,1fr); }}
  .visual.grid {{ grid-template-columns:repeat(2,1fr); }}
  .visual.six {{ grid-template-columns:repeat(3,1fr); }}
  .concept {{ border:2px solid var(--teal);border-radius:.08in;padding:.15in .16in;background:#F4FAF8;min-height:1.14in; }}
  .flow .concept,.compare .concept {{ min-height:2.0in; }}
  .concept h2 {{ margin:0 0 .09in;color:var(--teal);font-size:17pt;line-height:1.08; }}
  .concept p {{ margin:0;font-size:12pt;line-height:1.25; }}
  .concept.cobalt {{ border-color:var(--cobalt);background:#F2F6FF; }} .concept.cobalt h2 {{ color:var(--cobalt); }}
  .concept.maroon {{ border-color:var(--maroon);background:#FCF4F4; }} .concept.maroon h2 {{ color:var(--maroon); }}
  .concept.brown {{ border-color:var(--brown);background:#FBF8F2; }} .concept.brown h2 {{ color:var(--brown); }}
  .concept.amber {{ border-color:var(--amber);background:#FFF8EA; }} .concept.amber h2 {{ color:var(--amber); }}
  .concept.red {{ border-color:var(--red);background:#FFF4F2; }} .concept.red h2 {{ color:var(--red); }}
  .concept.green {{ border-color:var(--green);background:#F2FBF5; }} .concept.green h2 {{ color:var(--green); }}
  .callout {{ margin:.08in 0 0 .0;padding:.09in .15in;border:1px solid var(--teal);border-radius:.06in;background:#EAF3F0;color:var(--teal);font-size:12pt;font-weight:700;text-align:center; }}
  footer {{ display:flex;justify-content:space-between;gap:.25in;margin-top:.09in;color:var(--muted);font-size:8.5pt; }}
  footer a {{ color:var(--muted);text-decoration:none; }}
  @media screen {{ body {{ display:grid;gap:24px;padding:24px;justify-content:center; }} .slide {{ box-shadow:0 12px 36px rgba(0,0,0,.14); }} }}
  </style>
</head>
<body>
<main aria-label="Agent AI survey teaching deck">
{''.join(slide_html)}
</main>
</body>
</html>
'''
    path.write_text(document, encoding="utf-8")


def build_pdf(html_path: Path, pdf_path: Path) -> None:
    subprocess.run(
        [
            "weasyprint",
            "--pdf-variant",
            "pdf/ua-1",
            "--pdf-tags",
            "--custom-metadata",
            str(html_path),
            str(pdf_path),
        ],
        check=True,
    )


def build_contact_sheet(pdf_path: Path, output_path: Path) -> None:
    document = fitz.open(pdf_path)
    thumbs: list[Image.Image] = []
    for page in document:
        pixmap = page.get_pixmap(matrix=fitz.Matrix(0.65, 0.65), alpha=False)
        image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
        thumbs.append(image)
    document.close()

    columns = 5
    gap = 28
    label_height = 38
    thumb_width = 300
    thumb_height = round(thumb_width * 9 / 16)
    rows = math.ceil(len(thumbs) / columns)
    sheet = Image.new(
        "RGB",
        (
            columns * thumb_width + (columns + 1) * gap,
            rows * (thumb_height + label_height) + (rows + 1) * gap,
        ),
        rgb("bg"),
    )
    draw = ImageDraw.Draw(sheet)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial Bold.ttf", 18)
    except OSError:
        font = ImageFont.load_default()
    for index, image in enumerate(thumbs):
        image = image.resize((thumb_width, thumb_height), Image.Resampling.LANCZOS)
        row, column = divmod(index, columns)
        x = gap + column * (thumb_width + gap)
        y = gap + row * (thumb_height + label_height + gap)
        sheet.paste(image, (x, y))
        draw.text((x, y + thumb_height + 8), f"{index + 1:02d} · {SLIDES[index]['stage']}", fill=rgb("ink"), font=font)
    sheet.save(output_path, optimize=True)


def publish(paths: list[Path]) -> None:
    PUBLIC.mkdir(parents=True, exist_ok=True)
    for path in paths:
        shutil.copy2(path, PUBLIC / path.name)


def main() -> None:
    pptx_path = ROOT / f"{SLUG}.pptx"
    html_path = ROOT / f"{SLUG}.html"
    pdf_path = ROOT / f"{SLUG}.pdf"
    contact_path = ROOT / "contact-sheet.png"
    build_pptx(pptx_path)
    build_html(html_path)
    build_pdf(html_path, pdf_path)
    build_contact_sheet(pdf_path, contact_path)
    publish([pptx_path, html_path, pdf_path, contact_path])
    print(f"Built {len(SLIDES)} slides")
    for path in [pptx_path, html_path, pdf_path, contact_path]:
        print(path)
    print(f"Published to {PUBLIC}")


if __name__ == "__main__":
    main()
